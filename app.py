# app.py
import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import io
import tempfile

st.set_page_config(layout="wide")
st.title("📊 Purchase Order Report Generator")

uploaded_file = st.file_uploader("Upload your Excel work report", type=["xlsx"])

if uploaded_file:
    df = pd.read_excel(uploaded_file)
    df.columns = df.columns.str.strip()

    df["Creation Date"] = pd.to_datetime(df["Creation Date"], errors='coerce')
    df["Released Date"] = pd.to_datetime(df["Released Date"], errors='coerce')
    df["Days Difference"] = (df["Released Date"] - df["Creation Date"]).dt.days

    # Prepare metrics for dashboard
    total_pos = df["Purchase Order No."].nunique()
    approved_df = df[df["Overall Status"] == "APPROVED"]
    total_approved = approved_df["Purchase Order No."].nunique()
    in_progress = df[df["Overall Status"].str.upper().str.contains("PROGRESS", na=False)]
    total_in_progress = in_progress["Purchase Order No."].nunique()
    cancelled_deleted = df[df["Overall Status"].isin(["CANCELLED", "DELETED"])]
    total_cancelled = cancelled_deleted["Purchase Order No."].nunique()

    col_name = "Days Difference"
    df[col_name] = pd.to_numeric(df[col_name], errors="coerce")
    df = df[df[col_name] < 1000]
    valid_approvals = df[df[col_name].notnull()]
    average_days = round(valid_approvals[col_name].mean(), 2)
    delayed_approved = approved_df[approved_df[col_name] > 10]["Purchase Order No."].nunique()
    delayed_pct = round(100 * delayed_approved / total_approved, 2) if total_approved else 0
    waiting_pct = round(100 * total_in_progress / total_pos, 2) if total_pos else 0

    # Chart logic
    def dashboard_slide(ppt):
        slide_layout = ppt.slide_layouts[5]
        slide = ppt.slides.add_slide(slide_layout)
        tf = slide.shapes.title.text_frame
        tf.text = "Purchase Order workflow output"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True

        dashboard = [
            ("Total #POs", total_pos),
            ("Total Approved POs", total_approved),
            ("Total POs In Progress", total_in_progress),
            ("Total Cancelled/Deleted POs", total_cancelled),
            ("Average Approval Time/Approver", f"{average_days:.2f}"),
            ("Delayed Approvals > 10 days", delayed_approved),
            ("% Delayed Approval 'Approved PO'", f"{delayed_pct:.2f}%")
        ]

        rows = len(dashboard) + 1
        cols = 2
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(8.5), Inches(3)).table
        table.columns[0].width = Inches(5)
        table.columns[1].width = Inches(3)

        table.cell(0, 0).text = "Activity"
        table.cell(0, 1).text = "Value"

        for i, (label, value) in enumerate(dashboard):
            table.cell(i + 1, 0).text = label
            table.cell(i + 1, 1).text = str(value)

        # Callouts
        shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5), Inches(4), Inches(1))
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = RGBColor(199, 215, 238)
        shape1.text = f"Delaying Approvals >\n10 Days\n{delayed_pct:.2f}%"
        shape1.text_frame.paragraphs[0].font.size = Pt(18)
        shape1.text_frame.paragraphs[0].font.bold = True

        shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(5), Inches(4), Inches(1))
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = RGBColor(199, 215, 238)
        shape2.text = f"POs Waiting for\nApproval\n{waiting_pct:.2f}%"
        shape2.text_frame.paragraphs[0].font.size = Pt(18)
        shape2.text_frame.paragraphs[0].font.bold = True

    # Define chart options
    graph_options = {
        "📋 Dashboard Summary": None,
        "PO Avg Time by User: Approved": lambda d: d[d["Overall Status"] == "APPROVED"].groupby("Approver Name")["Days Difference"].mean(),
        "PO Count by User: Approved": lambda d: d[d["Overall Status"] == "APPROVED"].groupby("Approver Name")["Purchase Order No."].count(),
        "PO Count by User: In Progress": lambda d: d[d["Overall Status"] == "In Progress"].groupby("Approver Name")["Purchase Order No."].count(),
        "PO Avg Time by User: In Progress": lambda d: d[d["Overall Status"] == "In Progress"].groupby("Approver Name")["Days Difference"].mean(),
        "POs Cancelled/Deleted by Company": lambda d: d[d["Overall Status"].isin(["Cancelled", "DELETED"])].groupby("Company Code Decription")["Purchase Order No."].count(),
        "PO Avg Time by Company: Approved": lambda d: d[d["Overall Status"] == "APPROVED"].groupby("Company Code Decription")["Days Difference"].mean(),
        "PO Avg Time by Company: In Progress": lambda d: d[d["Overall Status"] == "In Progress"].groupby("Company Code Decription")["Days Difference"].mean(),
    }

    selected_graphs = st.multiselect("Select graphs to generate", list(graph_options.keys()))

    if selected_graphs:
        st.header("📈 Live Chart Previews")
        chart_images = []

        if "📋 Dashboard Summary" in selected_graphs:
            st.subheader("📋 Dashboard Summary")
            dashboard_data = {
                "Total #POs": total_pos,
                "Total Approved POs": total_approved,
                "Total POs In Progress": total_in_progress,
                "Total Cancelled/Deleted POs": total_cancelled,
                "Average Approval Time/Approver": average_days,
                "Delayed Approvals > 10 days": delayed_approved,
                "% Delayed Approval 'Approved PO'": f"{delayed_pct:.2f}%",
                "POs Waiting for Approval": f"{waiting_pct:.2f}%"
            }
            st.table(pd.DataFrame(dashboard_data.items(), columns=["Activity", "Value"]))

        for title in selected_graphs:
            if title == "📋 Dashboard Summary":
                continue
            st.subheader(title)
            chart_data = graph_options[title](df)
            if chart_data.empty:
                st.warning("No data for this chart.")
                continue

            fig, ax = plt.subplots(figsize=(6, 4))
            chart_data.sort_values().plot(kind="bar", ax=ax)
            ax.set_title(title)
            ax.set_ylabel("Value")
            ax.set_xticklabels(chart_data.index, rotation=45, ha='right')
            st.pyplot(fig)

            img_buf = io.BytesIO()
            fig.savefig(img_buf, format='png')
            chart_images.append((title, img_buf))
            plt.close(fig)

        if st.button("📥 Generate and Download PowerPoint"):
            ppt = Presentation()
            if "📋 Dashboard Summary" in selected_graphs:
                dashboard_slide(ppt)

            slide_layout = ppt.slide_layouts[5]
            for title, img_buf in chart_images:
                slide = ppt.slides.add_slide(slide_layout)
                slide.shapes.title.text = title
                img_buf.seek(0)
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
                    tmp_img.write(img_buf.read())
                    tmp_img_path = tmp_img.name
                slide.shapes.add_picture(tmp_img_path, Inches(1), Inches(1.5), width=Inches(8))

            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_ppt:
                ppt.save(tmp_ppt.name)
                st.success("✅ PowerPoint generated!")
                st.download_button("📤 Download .pptx", data=open(tmp_ppt.name, "rb"), file_name="PO_Report.pptx")

