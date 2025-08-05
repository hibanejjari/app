
import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
import tempfile

st.set_page_config(layout="wide")
st.title("📊 Purchase Order Report Generator")

uploaded_file = st.file_uploader("Upload your Excel work report", type=["xlsx"])

if uploaded_file:
    df = pd.read_excel(uploaded_file)
    df.columns = df.columns.str.strip()  # Clean column names

    # Ensure key columns exist
    if "Purchase Order No." in df.columns and "Serial No." in df.columns and "Overall Status" in df.columns:
        # Ensure correct types
        df["Purchase Order No."] = df["Purchase Order No."].astype(str).str.strip()
        df["Serial No."] = df["Serial No."].astype(int)

        # Create 'Latest Status' column based on max Serial No. per PO
        latest_rows = df.loc[df.groupby("Purchase Order No.")["Serial No."].idxmax()]
        po_to_status = latest_rows.set_index("Purchase Order No.")["Overall Status"]
        df["Latest Status"] = df["Purchase Order No."].map(po_to_status)

    # Clean relevant columns
    df["Overall Status"] = df["Overall Status"].astype(str).str.strip().str.upper()
    df["Approver Action"] = df["Approver Action"].astype(str).str.strip().str.upper()

    # Convert Days column
    col_name = "No.of Days to Approve"
    df[col_name] = pd.to_numeric(df[col_name].astype(str).str.replace(",", ""), errors="coerce")

    # Compute dashboard metrics
    total_pos = df["Purchase Order No."].nunique()
    approved_df = df[df["Latest Status"] == "APPROVED"]
    total_approved = approved_df["Purchase Order No."].nunique()
    in_progress = df["Overall Status"].str.contains("PROGRESS", na=False)
    total_in_progress = df[in_progress]["Purchase Order No."].nunique()
    cancelled_deleted = df["Overall Status"].isin(["CANCELLED", "DELETED"])
    total_cancelled = df[cancelled_deleted]["Purchase Order No."].nunique()

    avg_days_df = df[(df["Approver Action"] == "APPROVED") & (df[col_name].notnull())]
    average_days = round(avg_days_df[col_name].mean(), 2)

    delayed_rows = df[
        (df["Approver Action"] == "APPROVED") &
        (df[col_name] > 10) & (df[col_name] <= 20)
    ]
    delayed_approved = delayed_rows["Purchase Order No."].nunique()

    delayed_pct = round(100 * delayed_approved / total_approved, 2) if total_approved else 0
    waiting_pct = round(100 * total_in_progress / total_pos, 2) if total_pos else 0

    # --- Dashboard DataFrame for display ---
    dashboard_df = pd.DataFrame([
        ["Total #POs", total_pos],
        ["Total Approved POs", total_approved],
        ["Total POs In Progress", total_in_progress],
        ["Total Cancelled/Deleted POs", total_cancelled],
        ["Average Approval Time/Approver", average_days],
        ["Delayed Approvals > 10 days", delayed_approved],
        ["% Delayed Approval 'Approved PO'", f"{int(round(delayed_pct))}%"],
        ["POs Waiting for Approval", f"{int(round(waiting_pct))}%"],
    ], columns=["Activity", "Value"])

    # --- Chart logic ---
    graph_options = {
        "📋 Dashboard Summary": None,
        "PO Avg Time by User: Approved": lambda d: d[d["Latest Status"] == "APPROVED"].groupby("Approver Name")[col_name].mean(),
        "PO Count by User: Approved": lambda d: d[d["Latest Status"] == "APPROVED"].groupby("Approver Name")["Purchase Order No."].count(),
        "PO Count by User: In Progress": lambda d: d[d["Overall Status"] == "IN PROGRESS"].groupby("Approver Name")["Purchase Order No."].count(),
        "PO Avg Time by User: In Progress": lambda d: d[d["Overall Status"] == "IN PROGRESS"].groupby("Approver Name")[col_name].mean(),
        "POs Cancelled/Deleted by Company": lambda d: d[d["Overall Status"].isin(["CANCELLED", "DELETED"])].groupby("Company Code Decription")["Purchase Order No."].count(),
        "PO Avg Time by Company: Approved": lambda d: d[d["Latest Status"] == "APPROVED"].groupby("Company Code Decription")[col_name].mean(),
        "PO Avg Time by Company: In Progress": lambda d: d[d["Overall Status"] == "IN PROGRESS"].groupby("Company Code Decription")[col_name].mean(),
    }

    selected_graphs = st.multiselect("Select graphs to generate", list(graph_options.keys()), default=["📋 Dashboard Summary"])

    chart_images = []

    if selected_graphs:
        for title in selected_graphs:
            if title == "📋 Dashboard Summary":
                st.header("🧾 Dashboard Summary")
                st.dataframe(dashboard_df)

                # Save dashboard slide for PPT
                img_buf = io.BytesIO()
                # We skip image rendering of table; handled in PPT logic below
                chart_images.append(("📋 Dashboard Summary", "DASHBOARD_TABLE"))
                continue

            st.subheader(title)
            chart_data = graph_options[title](df)
            if chart_data.empty:
                st.warning("No data available for this chart.")
                continue

            fig, ax = plt.subplots(figsize=(14, 5))  # Wider, less zoomed-in
            sorted_data = chart_data.sort_values()
            bars = sorted_data.plot(kind="bar", ax=ax)
            ax.set_title(title, fontsize=14)
            ax.set_ylabel("Value", fontsize=12)
            ax.set_xticklabels(sorted_data.index, rotation=30, ha='right', fontsize=8)
            for p in bars.patches:
                value = round(p.get_height(), 2)
                ax.annotate(f'{value}', (p.get_x() + p.get_width() / 2, p.get_height()),
                ha='center', va='bottom', fontsize=8, xytext=(0, 3), textcoords='offset points')
            plt.tight_layout()
            st.pyplot(fig)


            # Save figure to buffer
            img_buf = io.BytesIO()
            fig.savefig(img_buf, format='png')
            chart_images.append((title, img_buf))
            plt.close(fig)

    # --- Generate PPT ---
    if st.button("📥 Generate and Download PowerPoint"):
        ppt = Presentation()
        slide_layout = ppt.slide_layouts[5]

        for title, img in chart_images:
            slide = ppt.slides.add_slide(slide_layout)
            slide.shapes.title.text = title

            if img == "DASHBOARD_TABLE":
                dashboard = dashboard_df.values.tolist()
                rows, cols = len(dashboard) + 1, 2
                table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(8.5), Inches(3)).table
                table.columns[0].width = Inches(5)
                table.columns[1].width = Inches(3)
                table.cell(0, 0).text = "Activity"
                table.cell(0, 1).text = "Value"

                for i, (label, value) in enumerate(dashboard):
                    table.cell(i + 1, 0).text = str(label)
                    table.cell(i + 1, 1).text = str(value)

            else:
                img.seek(0)
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
                    tmp_img.write(img.read())
                    slide.shapes.add_picture(tmp_img.name, Inches(1), Inches(1.5), width=Inches(8))

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_ppt:
            ppt.save(tmp_ppt.name)
            st.success("✅ PowerPoint generated!")
            st.download_button("📤 Download .pptx", data=open(tmp_ppt.name, "rb"), file_name="PO_Report.pptx")
